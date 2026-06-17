"""Text extraction for document formats other than PDF: docx, pptx, txt, md, csv.

Every format is reduced to plain text and chunked exactly like PDFs, so the
chunks flow through the same embedding → DocumentChunk → retrieval path. Only
pptx carries a meaningful page_number (the slide number); the rest use None.
"""
import csv as _csv

from ingestion.chunker import chunk_text
from config import get_settings

settings = get_settings()

# Extensions handled here (PDF and MP4 keep their own dedicated processors).
SUPPORTED_DOC_EXTENSIONS = {".docx", ".pptx", ".txt", ".md", ".csv"}


def _read_text_file(file_path: str) -> str:
    """Read a plain-text file, tolerating non-UTF-8 bytes."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
        return fh.read()


def _read_csv_file(file_path: str) -> str:
    """Flatten a CSV into readable 'col: value' lines so rows survive chunking."""
    lines: list[str] = []
    with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as fh:
        reader = _csv.reader(fh)
        rows = list(reader)
    if not rows:
        return ""
    header = rows[0]
    has_header = len(rows) > 1 and any(cell.strip() for cell in header)
    for row in (rows[1:] if has_header else rows):
        if has_header:
            cells = [f"{header[i].strip()}: {val.strip()}" for i, val in enumerate(row) if i < len(header)]
            lines.append(" | ".join(cells))
        else:
            lines.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(lines)


def _read_docx_file(file_path: str) -> str:
    """Extract paragraphs and table cells from a Word document, in order."""
    from docx import Document

    doc = Document(file_path)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx_segments(file_path: str) -> list[tuple[int, str]]:
    """Return (slide_number, text) per slide for a PowerPoint deck."""
    from pptx import Presentation

    prs = Presentation(file_path)
    segments: list[tuple[int, str]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            segments.append((slide_num, "\n".join(texts)))
    return segments


def extract_document_chunks(file_path: str, file_name: str, gcs_url: str, ext: str) -> list[dict]:
    """Dispatch by extension and return standard chunk dicts.

    `ext` is the lowercase file extension including the dot (e.g. ".docx").
    """
    source_type = ext.lstrip(".")
    all_chunks: list[dict] = []
    chunk_index = 0

    if ext == ".pptx":
        segments = _extract_pptx_segments(file_path)
    elif ext == ".docx":
        segments = [(None, _read_docx_file(file_path))]
    elif ext == ".csv":
        segments = [(None, _read_csv_file(file_path))]
    elif ext in {".txt", ".md"}:
        segments = [(None, _read_text_file(file_path))]
    else:
        raise ValueError(f"Extensión de documento no soportada: {ext}")

    for page_number, text in segments:
        if not text.strip():
            continue
        for chunk in chunk_text(text, settings.chunk_size, settings.chunk_overlap):
            all_chunks.append({
                "content": chunk,
                "source_type": source_type,
                "file_name": file_name,
                "gcs_url": gcs_url,
                "page_number": page_number,
                "chunk_index": chunk_index,
            })
            chunk_index += 1

    return all_chunks
